#!/usr/bin/env python3
"""
Generate a TimoRides pitch deck from the structured YAML schema.

Usage:
    python generate_pitch_deck.py \
        --schema ../timo_rides_new_deck.schema.yaml \
        --output output/TimoRides_Smart_Mobility_Intelligence.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Iterable, List, Tuple

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt


def _ensure_text_frame(
    text_frame,
    lines: Iterable[Tuple[str, int]],
) -> None:
    """Populate a text frame with lines and bullet levels."""
    paragraphs = list(lines)
    text_frame.clear()
    if not paragraphs:
        return

    for index, (text, level) in enumerate(paragraphs):
        if text is None:
            text = ""
        if not isinstance(text, str):
            text = str(text)
        para = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        para.text = text
        para.level = max(level, 0)


def _format_stats(stats: Iterable[dict]) -> List[Tuple[str, int]]:
    lines: List[Tuple[str, int]] = []
    for item in stats:
        label = item.get("label")
        value = item.get("value")
        if not value:
            continue
        headline = f"{label}: {value}" if label else value
        lines.append((headline, 0))
    return lines


def _format_named_descriptions(items: Iterable[dict]) -> List[Tuple[str, int]]:
    lines: List[Tuple[str, int]] = []
    for item in items:
        name = item.get("name") or item.get("heading")
        desc = item.get("description") or item.get("detail")
        if name:
            lines.append((name, 0))
        if desc:
            lines.append((desc, 1))
    return lines


def _generic_bullets(items: Iterable[str]) -> List[Tuple[str, int]]:
    return [(str(item), 0) for item in items if item]

def _add_slide(prs: Presentation, preferred_index: int = 1):
    """Add a slide using preferred layout index with graceful fallback."""
    total_layouts = len(prs.slide_layouts)
    index = preferred_index if preferred_index < total_layouts else 0
    return prs.slides.add_slide(prs.slide_layouts[index])


def _set_slide_title(slide, text: str | None) -> None:
    if text is None:
        return

    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        title_shape = slide.shapes.add_textbox(
            Inches(1.2),
            Inches(0.8),
            Inches(8.6),
            Inches(1.8),
        )
    tf = title_shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    para.text = text
    para.level = 0
    para.font.size = Pt(44)
    para.font.bold = True
    para.font.color.rgb = RGBColor(15, 56, 155)
    para.font.name = "Space Grotesk"


def _get_body_text_frame(slide):
    # search placeholders first
    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT):
                return shape.text_frame

    # fallback to first shape with text frame other than title
    title_shape = getattr(slide.shapes, "title", None)
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if shape.has_text_frame:
            return shape.text_frame

    # final resort: create textbox occupying central area
    textbox = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(2.2),
        Inches(8.6),
        Inches(4.8),
    )
    return textbox.text_frame


def _add_text_block(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Iterable[Tuple[str, int]],
    base_font_size: int = 24,
    bullet: bool = False,
) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = textbox.text_frame
    tf.clear()
    for idx, (text, level) in enumerate(lines):
        if text is None:
            continue
        if not isinstance(text, str):
            text = str(text)
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = text
        para.level = level if bullet else 0
        para.font.name = "Space Grotesk"
        para.font.size = Pt(max(base_font_size - (level * 2), 16))
        para.font.color.rgb = RGBColor(17, 24, 39)
        para.space_after = Pt(6)
        para.space_before = Pt(4)


def handle_cover(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 0)
    _set_slide_title(slide, slide_data.get("title", ""))

    subtitle_lines = []
    subtitle_text = slide_data.get("subtitle")
    if subtitle_text:
        subtitle_lines.append((subtitle_text, 0))
    if subtitle_lines:
        _add_text_block(
            slide,
            left=1.2,
            top=2.0,
            width=8.5,
            height=1.5,
            lines=subtitle_lines,
            base_font_size=28,
        )

    footer_lines = []
    footer = slide_data.get("footer")
    if footer:
        footer_lines.append((footer, 0))
        _add_text_block(
            slide,
            left=1.2,
            top=6.3,
            width=8.5,
            height=0.8,
            lines=footer_lines,
            base_font_size=16,
        )
    return slide


def handle_narrative(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title") or slide_data.get("headline", ""))
    text = slide_data.get("body") or slide_data.get("intro") or ""
    _add_text_block(
        slide,
        left=1.2,
        top=2.3,
        width=8.5,
        height=4.4,
        lines=[(text, 0)],
        base_font_size=24,
    )
    return slide


def handle_bullet_summary(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title") or slide_data.get("headline", ""))
    bullets = slide_data.get("bullets") or []
    formatted: List[Tuple[str, int]] = []
    for item in bullets:
        if isinstance(item, dict):
            for key, value in item.items():
                formatted.append((f"{key}: {value}", 0))
        else:
            formatted.append((str(item), 0))
    _add_text_block(
        slide,
        left=1.2,
        top=2.4,
        width=8.5,
        height=4.5,
        lines=formatted,
        base_font_size=22,
        bullet=True,
    )
    return slide


def handle_market_beachhead(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    stats = slide_data.get("stats", [])
    _add_text_block(
        slide,
        left=1.2,
        top=2.4,
        width=8.5,
        height=4.5,
        lines=_format_stats(stats),
        base_font_size=24,
    )
    return slide


def handle_differentiators(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    pillars = slide_data.get("pillars", [])
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=4.1,
        height=4.6,
        lines=_format_named_descriptions(pillars[:len(pillars)//2 + len(pillars)%2]),
        base_font_size=20,
    )
    if len(pillars) > 1:
        _add_text_block(
            slide,
            left=5.6,
            top=2.2,
            width=4.1,
            height=4.6,
            lines=_format_named_descriptions(pillars[len(pillars)//2 + len(pillars)%2:]),
            base_font_size=20,
        )
    return slide


def handle_integrations(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    body_text = slide_data.get("body")
    if body_text:
        lines.append((body_text, 0))
    lines.extend(_format_named_descriptions(slide_data.get("features", [])))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_problem_landscape(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    if slide_data.get("intro"):
        lines.append((slide_data["intro"], 0))
    for segment in slide_data.get("segments", []):
        audience = segment.get("audience")
        if audience:
            lines.append((audience, 0))
        for challenge in segment.get("challenges", []):
            heading = challenge.get("heading")
            if heading:
                lines.append((heading, 1))
            for bullet in challenge.get("bullets", []):
                lines.append((bullet, 2))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=21,
    )
    return slide


def handle_solution(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    for segment in slide_data.get("segments", []):
        heading = segment.get("heading")
        if heading:
            lines.append((heading, 0))
        for bullet in segment.get("bullets", []):
            lines.append((bullet, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_product_flow(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    if slide_data.get("caption"):
        lines.append((slide_data["caption"], 0))
    for step in slide_data.get("steps", []):
        label = step.get("label")
        if label:
            lines.append((label, 0))
        desc = step.get("description")
        if desc:
            lines.append((desc, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_comparison(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    for point in slide_data.get("comparison_points", []):
        attribute = point.get("attribute")
        if attribute:
            lines.append((attribute, 0))
        timo = point.get("timo")
        incumbents = point.get("incumbents")
        why = point.get("why_it_matters")
        if timo:
            lines.append((f"Timo: {timo}", 1))
        if incumbents:
            lines.append((f"Incumbents: {incumbents}", 1))
        if why:
            lines.append((f"Why it matters: {why}", 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=21,
    )
    return slide


def handle_metrics(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    metrics = slide_data.get("metrics", [])
    _add_text_block(
        slide,
        left=1.2,
        top=2.3,
        width=8.5,
        height=4.5,
        lines=_format_stats(metrics),
        base_font_size=24,
    )
    return slide


def handle_growth_levers(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    if slide_data.get("overview"):
        lines.append((slide_data["overview"], 0))
    for target in slide_data.get("targets", []):
        name = target.get("name")
        baseline = target.get("baseline")
        goal = target.get("goal")
        if name:
            lines.append((name, 0))
        if baseline or goal:
            lines.append((f"Baseline: {baseline} → Goal: {goal}", 1))
    if slide_data.get("levers"):
        lines.append(("Levers", 0))
        for lever in slide_data["levers"]:
            lines.append((lever, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.3,
        width=8.5,
        height=4.5,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_simple_points(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    points = slide_data.get("points") or slide_data.get("bullets") or []
    _add_text_block(
        slide,
        left=1.2,
        top=2.4,
        width=8.5,
        height=4.5,
        lines=_generic_bullets(points),
        base_font_size=22,
        bullet=True,
    )
    return slide


def handle_market_opportunity(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    stats = slide_data.get("stats", [])
    _add_text_block(
        slide,
        left=1.2,
        top=2.3,
        width=8.5,
        height=4.5,
        lines=_format_stats(stats),
        base_font_size=24,
    )
    return slide


def handle_team(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    entries = slide_data.get("team", [])
    lines: List[Tuple[str, int]] = []
    for person in entries:
        name = person.get("name")
        role = person.get("role")
        blurb = person.get("blurb")
        link = person.get("link")
        if name:
            title_line = f"{name} – {role}" if role else name
            lines.append((title_line, 0))
        if blurb:
            lines.append((blurb, 1))
        if link:
            lines.append((link, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=21,
    )
    return slide


def handle_advisors(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    entries = slide_data.get("advisors", [])
    lines: List[Tuple[str, int]] = []
    for advisor in entries:
        name = advisor.get("name")
        role = advisor.get("role")
        blurb = advisor.get("blurb")
        if name:
            role_line = f"{name} – {role}" if role else name
            lines.append((role_line, 0))
        if blurb:
            lines.append((blurb, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=21,
    )
    return slide


def handle_traction(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    highlights = slide_data.get("highlights", [])
    lines: List[Tuple[str, int]] = []
    for item in highlights:
        metric = item.get("metric")
        desc = item.get("description")
        if metric:
            lines.append((metric, 0))
        if desc:
            lines.append((desc, 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.3,
        width=8.5,
        height=4.5,
        lines=lines,
        base_font_size=24,
    )
    return slide


def handle_fundraising(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    if slide_data.get("raise_amount"):
        lines.append((f"Raise amount: {slide_data['raise_amount']}", 0))
    if slide_data.get("valuation_cap"):
        lines.append((f"Valuation: {slide_data['valuation_cap']}", 0))
    if slide_data.get("discount"):
        lines.append((f"Discount: {slide_data['discount']}", 0))
    if slide_data.get("fund_allocation"):
        lines.append(("Use of funds", 0))
        for item in slide_data["fund_allocation"]:
            label = item.get("label")
            pct = item.get("percentage")
            if label:
                lines.append((f"{label} – {pct}", 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_revenue_model(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []

    fare_model = slide_data.get("fare_model", {})
    if fare_model:
        lines.append(("Fare structure", 0))
        for key in ("base_fare", "price_per_km", "price_per_minute", "minimum_fare"):
            if fare_model.get(key) is not None:
                label = key.replace("_", " ").title()
                lines.append((f"{label}: {fare_model[key]}", 1))

    avg = slide_data.get("average_ride_assumptions")
    if avg:
        lines.append(("Average ride assumptions", 0))
        if avg.get("distance_km") is not None:
            lines.append((f"Distance: {avg['distance_km']} km", 1))
        if avg.get("duration_minutes") is not None:
            lines.append((f"Duration: {avg['duration_minutes']} minutes", 1))

    profitability = slide_data.get("profitability", {})
    if profitability:
        lines.append(("Profitability metrics", 0))
        for key, value in profitability.items():
            label = key.replace("_", " ").title()
            lines.append((f"{label}: {value}", 1))

    segments = slide_data.get("segments")
    if segments:
        lines.append(("Customer segments", 0))
        for segment in segments:
            lines.append((segment, 1))

    strategy = slide_data.get("strategy")
    if strategy:
        lines.append(("Growth strategy", 0))
        lines.append((strategy, 1))

    risks = slide_data.get("risks_and_mitigation")
    if risks:
        lines.append(("Risks & mitigation", 0))
        for item in risks:
            risk = item.get("risk")
            mitigation = item.get("mitigation")
            if risk:
                lines.append((f"Risk: {risk}", 1))
            if mitigation:
                lines.append((f"Mitigation: {mitigation}", 2))

    _add_text_block(
        slide,
        left=1.2,
        top=2.2,
        width=8.5,
        height=4.6,
        lines=lines,
        base_font_size=20,
    )
    return slide


def handle_closing(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    _set_slide_title(slide, slide_data.get("title"))
    lines: List[Tuple[str, int]] = []
    if slide_data.get("body"):
        lines.append((slide_data["body"], 0))
    if slide_data.get("footer"):
        lines.append((slide_data["footer"], 1))
    _add_text_block(
        slide,
        left=1.2,
        top=2.4,
        width=8.5,
        height=4.4,
        lines=lines,
        base_font_size=22,
    )
    return slide


def handle_generic(prs: Presentation, slide_data: dict):
    slide = _add_slide(prs, 1)
    title = slide_data.get("title") or slide_data.get("headline") or f"Slide {slide_data.get('id', '')}"
    _set_slide_title(slide, title)
    raw_texts = []
    for key, value in slide_data.items():
        if key in {"id", "layout", "title", "headline"}:
            continue
        if isinstance(value, str):
            raw_texts.append(value)
    _add_text_block(
        slide,
        left=1.2,
        top=2.4,
        width=8.5,
        height=4.4,
        lines=_generic_bullets(raw_texts),
        base_font_size=22,
        bullet=True,
    )
    return slide


LAYOUT_HANDLERS = {
    "cover": handle_cover,
    "narrative": handle_narrative,
    "mission_perspective": handle_narrative,
    "closing": handle_closing,
    "bullet_summary": handle_bullet_summary,
    "market_beachhead": handle_market_beachhead,
    "differentiators_grid": handle_differentiators,
    "integrations": handle_integrations,
    "problem_landscape": handle_problem_landscape,
    "solution_response": handle_solution,
    "product_flow": handle_product_flow,
    "differentiators_comparison": handle_comparison,
    "unit_economics_snapshot": handle_metrics,
    "mission_targets": handle_bullet_summary,
    "growth_levers": handle_growth_levers,
    "risk_mitigation": handle_simple_points,
    "market_opportunity": handle_market_opportunity,
    "team_core": handle_team,
    "team_extended": handle_team,
    "advisors": handle_advisors,
    "mission_statement": handle_bullet_summary,
    "traction": handle_traction,
    "fundraising": handle_fundraising,
    "revenue_model": handle_revenue_model,
}


def parse_args(argv: List[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a pitch deck from schema.")
    parser.add_argument(
        "--schema",
        required=True,
        type=Path,
        help="Path to the YAML schema file.",
    )
    parser.add_argument(
        "--output",
        required=True,
        type=Path,
        help="Where to write the generated .pptx.",
    )
    parser.add_argument(
        "--template",
        type=Path,
        help="Optional .pptx template to use as a base.",
    )
    return parser.parse_args(argv)


def load_schema(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def build_presentation(schema: dict, template: Path | None) -> Presentation:
    prs = Presentation(str(template)) if template else Presentation()
    # remove any existing slides from the template to avoid leftover content
    if template and len(prs.slides) > 0:
        slide_ids = list(prs.slides._sldIdLst)  # pylint: disable=protected-access
        for slide_id in slide_ids:
            r_id = slide_id.rId  # type: ignore[attr-defined]
            prs.part.drop_rel(r_id)
            prs.slides._sldIdLst.remove(slide_id)  # pylint: disable=protected-access
    slides = schema.get("slides", [])
    for slide_data in slides:
        layout = slide_data.get("layout")
        handler = LAYOUT_HANDLERS.get(layout, handle_generic)
        handler(prs, slide_data)
    return prs


def main(argv: List[str]) -> int:
    args = parse_args(argv)
    schema = load_schema(args.schema)

    prs = build_presentation(schema, args.template)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"Deck written to {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

